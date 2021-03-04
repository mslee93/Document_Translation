from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from . import papa
import re

check_word = re.compile(r'[`~!@#$£€¥%^&*\(\)-_+=\\\|\]\}\[\{\"\';:/?.>,<＃＆＊＠§※☆★○●◎◇◆□■△▲▽▼→←↑↓↔〓◁◀▷▶♤♠♡'
                        r'♥♧♣⊙◈▣◐◑☏☎¶†‡↕↗↙↖↘♬㈜®ºª＂（）［］｛｝‘’“”〔〕〈〉《》「」『』【】ⓐⓑⓒⓓⓔⓕⓖⓗⓘⓙⓚ'
                        r'ⓛ①②③④⑤⑥⑦⑧⑨⒜⒝⒞⒟⒠⒡⒢⒣⒤⒥￦Ｆ′″℃¤℉½⅓⅔¼¼ⅰⅱⅲⅳⅴⅵⅶⅠⅡⅢⅣⅤⅥⅦⅧⅨⅩ＋－＜＝＞±×÷≠≤≥∞∴'
                        r'♂♀∇≡≒≪≫√∽∝∵∫∬∈∋⊆⊇⊂⊃∪∩∧∨⇒￢⇔∮∑∏　！＇，．／：；？＾＿｀｜￣、。·‥…¨〃­­­­∥＼∼´¡'
                        r'―¿ㅇˇ˘˝˚˙¸˛¡ː㉠㉡㉢㉣㉤㉥㉦㉧㉨㉩㉮㉯㉰㉱㉲㉳㉴㈀㈁㈂㈃㈄㈅㈆㈎㈏㈐㈑㈒㈓㈔㈕\d\s]*')

engine_dic = {'papago': papa}
engine = None

# pptx의 text를 속성 그대로 살려두고 글자만 변경
def replace_paragraph_text_retaining_initial_formatting(paragraph, new_text):
    if len(paragraph.runs) == 0:
        return
    p = paragraph._p
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)
    paragraph.runs[0].text = new_text


def translate(shape, src_lang, tgt_lang):
    for paragraph in shape.text_frame.paragraphs:
        # 문자열 예외 처리 공백..
        if paragraph.text.isspace() or paragraph.text is '':
            continue
        if check_word.fullmatch(paragraph.text):
            continue
        # 문자열이 번역할 언어인지 확인
#       if papa.check_language(paragraph.text) == src_lang:
        # tgt_lang 언어로 번역
        new_text = engine.translate(paragraph.text, src_lang, tgt_lang)
        # 리턴 값이 none이라면 API 에러로 더이상 진행이 안되므로 바로 리턴
        if new_text is False:
            continue
        # 문자열 표시
        print('src text :' + paragraph.text + ", target text : " + new_text)
        # pptx 해당 문자열 변경
        replace_paragraph_text_retaining_initial_formatting(paragraph, new_text)


# 테이블 안에 있는 문자열 번역
def table_func(shape, src_lang, tgt_lang):
    table = shape.table
    # 셀단위로 체크
    for cell in table.iter_cells():
        translate(cell, src_lang, tgt_lang)


# 텍스트 상자 문자열 번역
def text_func(shape, src_lang, tgt_lang):
    translate(shape, src_lang, tgt_lang)


def chart_func(shape, src_lang, tgt_lang):
    chart = shape.chart
    if chart.has_title:
        if chart.chart_title.has_text_frame:
            translate(chart.chart_title, src_lang, tgt_lang)
    try:
        if chart.value_axis.has_title:
            if chart.value_axis.axis_title.has_text_frame:
                translate(chart.value_axis.axis_title, src_lang, tgt_lang)
    except: 
        pass
    try:
        if chart.category_axis.has_title:
            if chart.category_axis.axis_title.has_text_frame:
                translate(chart.category_axis.axis_title, src_lang, tgt_lang)
    except:
        pass
    # if chart.plots[0].has_data_labels:
    #     print(5)
    #     if chart.plots[0].data_labels.has_text_frame:
    #         print(6)
    #         translate(chart.plots[0].data_labels, src_lang, tgt_lang)

    # if chart.series.BarSeries:
    #     print(7)
    #     if chart.plots.data_labels.has_text_frame:
    #         print(8)
    #         translate(chart.plots.data_labels, src_lang, tgt_lang)


def handle_shape(prs, tgt_file, shape, src_lang, tgt_lang):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for _shape in shape.shapes:
            handle_shape(prs, tgt_file, _shape, src_lang, tgt_lang)
    # 객체가 테이블이면
    elif shape.has_table:
        # 테이블 번역 함수 호출
        table_func(shape, src_lang, tgt_lang)
            # 번역함수에서 리턴이 False라면 웹 API에 이상으로 지금까지 작업한 내용 저장하고 리턴
            # prs.save(tgt_file)
            # return False
    # 텍스트 상자면
    elif shape.has_text_frame:
        # 텍스트 번역 함수 호출
        text_func(shape, src_lang, tgt_lang)
            # 번역함수에서 리턴이 False라면 웹 API에 이상으로 지금까지 작업한 내용 저장하고 리턴
            # prs.save(tgt_file)
            # return False
    elif shape.has_chart:
        chart_func(shape, src_lang, tgt_lang)


def run(src_file, tgt_file, src_lang, tgt_lang, eng):
    global engine
    engine = engine_dic[eng]
    # pptx 불러오기
    prs = Presentation(src_file)
    # 슬라이드 단위로 루프
    for slide in prs.slides:
        # 슬라이드에서 각 객체(모양?) 단위로 루프
        for shape in slide.shapes:
            handle_shape(prs, tgt_file, shape, src_lang, tgt_lang)
    # 끝까지 문제없이 번역되었다면 저장하고 리턴
    prs.save(tgt_file)


if __name__ == "__main__":
    src_file = 'test_test3.pptx'
    tgt_file = 'output_test3_ch.pptx'
    src_lang = 'ko'
    tgt_lang = 'zh-CN'
    eng = 'papago'

    run(src_file, tgt_file, src_lang, tgt_lang, eng)
